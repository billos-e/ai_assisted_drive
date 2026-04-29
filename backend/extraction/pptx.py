from __future__ import annotations

from io import BytesIO

from pptx import Presentation


def extract_pptx_bytes(content: bytes) -> str:
    presentation = Presentation(BytesIO(content))
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                texts.append(text)
    return "\n".join(texts)
